"""
授業スライド自動生成ツール（歴史総合・公共）
Anthropic Claude を使ったインタビュー型スライド生成
"""

import os
import json
import re
from pathlib import Path
from datetime import datetime
from docx import Document
import anthropic
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO
from dotenv import load_dotenv

load_dotenv()
client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
MODEL = "claude-sonnet-4-6"

BASE_DIR = Path(__file__).parent
INPUT_TEXTS = BASE_DIR / "input" / "texts"
INPUT_MEMOS = BASE_DIR / "input" / "memos"
OUTPUT_DIR = BASE_DIR / "output"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ─── カラーパレット ───────────────────────────────
C = {
    "navy":        RGBColor(0x0D, 0x21, 0x37),
    "navy2":       RGBColor(0x1E, 0x4D, 0x7B),
    "orange":      RGBColor(0xE8, 0x60, 0x1A),
    "orange_dk":   RGBColor(0xC0, 0x45, 0x08),
    "teal":        RGBColor(0x1B, 0x8B, 0x6F),
    "teal_dk":     RGBColor(0x0F, 0x5C, 0x47),
    "white":       RGBColor(0xFF, 0xFF, 0xFF),
    "off_white":   RGBColor(0xF4, 0xF6, 0xF9),
    "light_bg":    RGBColor(0xEB, 0xEF, 0xF5),
    "dark":        RGBColor(0x12, 0x16, 0x24),
    "mid":         RGBColor(0x4A, 0x5B, 0x72),
    "cream":       RGBColor(0xFF, 0xF5, 0xEB),
}


# ─── 入力ファイル読み込み ────────────────────────────

def read_docx(path: Path) -> str:
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def load_dir(directory: Path) -> str:
    texts = []
    for path in sorted(directory.glob("*")):
        if path.suffix == ".txt":
            texts.append(path.read_text(encoding="utf-8"))
        elif path.suffix == ".docx":
            texts.append(read_docx(path))
    return "\n\n".join(texts)


# ─── シェイプ描画ユーティリティ ──────────────────────

def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(MSO.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_oval(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(MSO.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text: str, left, top, width, height,
             size=26, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Meiryo UI"
    if color:
        run.font.color.rgb = color
    return txb


def add_numbered_item(slide, number: int, text: str, y_pos,
                      badge_color: RGBColor, text_color: RGBColor):
    """丸バッジ＋テキストの箇条書きアイテム"""
    badge_size = Inches(0.52)
    bx = Inches(0.62)
    by = y_pos + Inches(0.06)
    add_oval(slide, bx, by, badge_size, badge_size, badge_color)
    add_text(slide, str(number), bx, by + Inches(0.04),
             badge_size, badge_size - Inches(0.08),
             size=17, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
    add_text(slide, text, Inches(1.35), y_pos,
             SLIDE_W - Inches(1.95), Inches(0.72),
             size=25, color=text_color)


# ─── スライド種別ビルダー ────────────────────────────

def build_title_slide(prs: Presentation, data: dict, lesson_title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C["navy"])
    W, H = SLIDE_W, SLIDE_H

    # 右下装飾円（2層）
    add_oval(slide, W - Inches(4.2), H - Inches(4.2), Inches(6.5), Inches(6.5), C["navy2"])
    add_oval(slide, W - Inches(2.6), H - Inches(2.6), Inches(4.0), Inches(4.0), C["navy"])

    # 左オレンジ縦バー（グラデーション風に2本）
    add_rect(slide, 0, 0, Inches(0.28), H, C["orange_dk"])
    add_rect(slide, Inches(0.28), 0, Inches(0.08), H, C["orange"])

    # 科目ラベル
    add_rect(slide, Inches(0.55), Inches(0.45), Inches(1.7), Inches(0.48), C["orange"])
    add_text(slide, "公　共",
             Inches(0.55), Inches(0.45), Inches(1.7), Inches(0.48),
             size=18, bold=True, color=C["white"], align=PP_ALIGN.CENTER)

    # メインタイトル
    add_text(slide, lesson_title,
             Inches(0.55), Inches(1.7), W - Inches(1.3), Inches(2.8),
             size=50, bold=True, color=C["white"])

    # オレンジ区切り線
    add_rect(slide, Inches(0.55), Inches(4.6), Inches(4.8), Inches(0.07), C["orange"])

    # サブタイトル
    if data.get("body"):
        add_text(slide, data["body"],
                 Inches(0.55), Inches(4.82), W - Inches(1.3), Inches(0.9),
                 size=26, color=RGBColor(0xA8, 0xC4, 0xE0), italic=True)

    # 右下クレジット
    add_text(slide, "高校　公共",
             W - Inches(3.8), H - Inches(0.55), Inches(3.5), Inches(0.4),
             size=14, color=RGBColor(0x3A, 0x5A, 0x80), align=PP_ALIGN.RIGHT)


def build_organizer_slide(prs: Presentation, data: dict):
    """先行オーガナイザー（濃紺背景・ロードマップ風）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C["navy"])
    W, H = SLIDE_W, SLIDE_H

    # 右上装飾円
    add_oval(slide, W - Inches(3.5), -Inches(1.5), Inches(5.0), Inches(5.0), C["navy2"])

    # 左オレンジバー
    add_rect(slide, 0, 0, Inches(0.28), H, C["orange_dk"])
    add_rect(slide, Inches(0.28), 0, Inches(0.08), H, C["orange"])

    # ラベル
    add_rect(slide, Inches(0.55), Inches(0.38), Inches(4.5), Inches(0.52), C["orange"])
    add_text(slide, "  この授業のロードマップ",
             Inches(0.55), Inches(0.38), Inches(4.5), Inches(0.52),
             size=19, bold=True, color=C["white"])

    # タイトル
    add_text(slide, data.get("title", ""),
             Inches(0.55), Inches(1.05), W - Inches(0.9), Inches(0.9),
             size=30, bold=True, color=C["white"])

    # セパレーター
    add_rect(slide, Inches(0.55), Inches(2.08), W - Inches(0.9), Inches(0.04),
             RGBColor(0x3A, 0x5A, 0x80))

    bullets = data.get("bullets", [])
    n = len(bullets)
    item_h = Inches(1.05) if n >= 4 else Inches(1.2)
    y_start = Inches(2.28)
    for i, item in enumerate(bullets[:5]):
        y = y_start + i * item_h
        badge_size = Inches(0.52)
        add_oval(slide, Inches(0.62), y + Inches(0.06),
                 badge_size, badge_size, C["orange"])
        add_text(slide, str(i + 1),
                 Inches(0.62), y + Inches(0.1), badge_size, badge_size - Inches(0.08),
                 size=17, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        add_text(slide, item, Inches(1.35), y, W - Inches(1.95), Inches(0.82),
                 size=25, color=RGBColor(0xD0, 0xE4, 0xF8))


def build_content_slide(prs: Presentation, data: dict, slide_num: int = 0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C["off_white"])
    W, H = SLIDE_W, SLIDE_H

    # 右下装飾円（薄いライトブルー）
    add_oval(slide, W - Inches(3.8), H - Inches(3.8),
             Inches(5.5), Inches(5.5), C["light_bg"])

    # ヘッダー帯（ネイビー）
    add_rect(slide, 0, 0, W, Inches(1.35), C["navy"])
    add_rect(slide, 0, 0, Inches(0.28), Inches(1.35), C["orange_dk"])
    add_rect(slide, Inches(0.28), 0, Inches(0.08), Inches(1.35), C["orange"])

    add_text(slide, data.get("title", ""),
             Inches(0.55), Inches(0.22), W - Inches(0.75), Inches(0.92),
             size=30, bold=True, color=C["white"])

    # セパレーター
    add_rect(slide, Inches(0.55), Inches(1.5), W - Inches(0.75), Inches(0.04),
             RGBColor(0xCC, 0xD5, 0xE0))

    bullets = data.get("bullets", [])
    if bullets:
        n = len(bullets)
        item_h = Inches(1.25) if n <= 3 else Inches(1.05)
        y_start = Inches(1.68)
        for i, item in enumerate(bullets[:5]):
            add_numbered_item(slide, i + 1, item,
                              y_start + i * item_h, C["orange"], C["dark"])
    elif data.get("body"):
        add_text(slide, data["body"],
                 Inches(0.7), Inches(1.65), W - Inches(1.4), H - Inches(2.2),
                 size=26, color=C["dark"])

    # 下部ラインとスライド番号
    add_rect(slide, 0, H - Inches(0.1), W, Inches(0.1), C["orange"])
    if slide_num:
        add_text(slide, str(slide_num),
                 W - Inches(0.65), H - Inches(0.48),
                 Inches(0.5), Inches(0.38),
                 size=14, color=C["mid"], align=PP_ALIGN.RIGHT)


def build_question_slide(prs: Presentation, data: dict):
    """フルオレンジ背景で劇的な問いスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C["orange"])
    W, H = SLIDE_W, SLIDE_H

    # 右下に濃いオレンジ装飾円（2層）
    add_oval(slide, W - Inches(5.0), H - Inches(5.0),
             Inches(7.2), Inches(7.2), C["orange_dk"])
    add_oval(slide, W - Inches(2.0), H - Inches(2.0),
             Inches(3.5), Inches(3.5), C["orange"])

    # 左白縦バー
    add_rect(slide, 0, 0, Inches(0.28), H, C["white"])

    # 上部ラベル帯
    add_rect(slide, Inches(0.48), Inches(0.36), Inches(3.4), Inches(0.52), C["orange_dk"])
    add_text(slide, "  考えてみよう",
             Inches(0.48), Inches(0.36), Inches(3.4), Inches(0.52),
             size=20, bold=True, color=C["white"])

    # 白区切り線
    add_rect(slide, Inches(0.48), Inches(1.05), W - Inches(0.96), Inches(0.05), C["white"])

    # 問いテキスト（大・白）
    add_text(slide, data.get("title", ""),
             Inches(0.6), Inches(1.28), W - Inches(1.2), Inches(3.8),
             size=38, bold=True, color=C["white"])

    # 補足テキスト
    if data.get("body"):
        add_rect(slide, Inches(0.6), Inches(5.3), Inches(3.8), Inches(0.05),
                 RGBColor(0xFF, 0xD0, 0xA0))
        add_text(slide, data["body"],
                 Inches(0.6), Inches(5.45), W - Inches(1.2), Inches(1.75),
                 size=22, color=C["cream"])


def build_activity_slide(prs: Presentation, data: dict):
    """活動スライド（ミントグリーン背景）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RGBColor(0xF0, 0xFA, 0xF5))
    W, H = SLIDE_W, SLIDE_H

    # 右下装飾円
    add_oval(slide, W - Inches(3.5), H - Inches(3.5),
             Inches(5.2), Inches(5.2), RGBColor(0xD0, 0xF0, 0xE8))

    # ヘッダー（ティール）
    add_rect(slide, 0, 0, W, Inches(1.35), C["teal"])
    add_rect(slide, 0, 0, Inches(0.28), Inches(1.35), C["teal_dk"])
    add_rect(slide, Inches(0.28), 0, Inches(0.08), Inches(1.35), C["teal"])

    add_text(slide, data.get("title", ""),
             Inches(0.55), Inches(0.22), W - Inches(0.75), Inches(0.92),
             size=30, bold=True, color=C["white"])

    if data.get("body"):
        add_rect(slide, Inches(0.7), Inches(1.65), W - Inches(1.4), H - Inches(2.3),
                 RGBColor(0xDC, 0xF5, 0xEC))
        add_text(slide, data["body"],
                 Inches(1.0), Inches(1.95), W - Inches(2.0), H - Inches(3.0),
                 size=26, color=C["dark"])

    add_rect(slide, 0, H - Inches(0.1), W, Inches(0.1), C["teal"])


def build_summary_slide(prs: Presentation, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C["off_white"])
    W, H = SLIDE_W, SLIDE_H

    # 右下装飾円（ティール）
    add_oval(slide, W - Inches(3.5), H - Inches(3.5),
             Inches(5.2), Inches(5.2), RGBColor(0xD0, 0xF0, 0xE8))

    # ヘッダー（ティール）
    add_rect(slide, 0, 0, W, Inches(1.35), C["teal"])
    add_rect(slide, 0, 0, Inches(0.28), Inches(1.35), C["teal_dk"])
    add_rect(slide, Inches(0.28), 0, Inches(0.08), Inches(1.35), C["teal"])

    add_text(slide, data.get("title", "まとめ"),
             Inches(0.55), Inches(0.22), W - Inches(0.75), Inches(0.92),
             size=30, bold=True, color=C["white"])

    add_rect(slide, Inches(0.55), Inches(1.5), W - Inches(0.75), Inches(0.04),
             RGBColor(0xA0, 0xD4, 0xC0))

    bullets = data.get("bullets", [])
    n = len(bullets)
    item_h = Inches(1.25) if n <= 3 else Inches(1.05)
    y_start = Inches(1.68)
    for i, item in enumerate(bullets[:5]):
        y = y_start + i * item_h
        badge_size = Inches(0.52)
        add_oval(slide, Inches(0.62), y + Inches(0.06),
                 badge_size, badge_size, C["teal"])
        add_text(slide, str(i + 1),
                 Inches(0.62), y + Inches(0.1), badge_size, badge_size - Inches(0.08),
                 size=17, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        add_text(slide, item, Inches(1.35), y, W - Inches(1.95), Inches(0.82),
                 size=25, color=C["dark"])

    add_rect(slide, 0, H - Inches(0.1), W, Inches(0.1), C["teal"])


# ─── PPTX 組み立て ──────────────────────────────

def create_pptx(structure: dict, output_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    lesson_title = structure.get("title", "授業スライド")
    content_count = 0

    for slide_data in structure.get("slides", []):
        stype = slide_data.get("type", "content")

        if stype == "title":
            build_title_slide(prs, slide_data, lesson_title)
        elif stype == "organizer":
            build_organizer_slide(prs, slide_data)
        elif stype == "question":
            build_question_slide(prs, slide_data)
        elif stype == "activity":
            build_activity_slide(prs, slide_data)
        elif stype == "summary":
            build_summary_slide(prs, slide_data)
        else:
            content_count += 1
            build_content_slide(prs, slide_data, slide_num=content_count)

    prs.save(output_path)
    print(f"\n[完了] スライドを保存しました: {output_path}")
    print(f"[枚数] {len(structure.get('slides', []))}枚")


# ─── フェーズ1：インタビュー ─────────────────────────

INTERVIEW_SYSTEM = """\
あなたは高校「公共」「歴史総合」の授業設計をサポートするAIアシスタントです。
教師が授業スライドを作るにあたり、3〜5個の質問を通じて以下を明らかにします。

- 授業で最も伝えたいこと（核心メッセージ）
- 生徒に深く考えさせたい問いや論点
- 授業の展開スタイル（講義型・発問中心・資料読解など）
- 生徒の実態や前提知識

【ルール】
- 1回に1つだけ質問する（50字以内）
- すべての質問が終わったら「ありがとうございました。」で締めくくる
"""


def run_interview(title: str, textbook: str, memo: str) -> list[dict]:
    context_msg = f"""
授業タイトル：{title}

【教科書テキスト】
{textbook if textbook else "（未提供）"}

【NotebookLMまとめ・メモ】
{memo if memo else "（未提供）"}

上記の授業について3〜5個の質問で方向性を探ってください。
まず最初の質問を1つだけしてください。
""".strip()

    messages = [{"role": "user", "content": context_msg}]

    print("\n" + "=" * 50)
    print("【インタビュー開始】")
    print("（終了したい場合は「終了」と入力）")
    print("=" * 50 + "\n")

    for _ in range(5):
        response = client.messages.create(
            model=MODEL, max_tokens=512,
            system=INTERVIEW_SYSTEM, messages=messages,
        )
        ai_reply = response.content[0].text
        messages.append({"role": "assistant", "content": ai_reply})
        print(f"AI: {ai_reply}\n")

        if "ありがとうございました" in ai_reply:
            break

        user_input = input("あなた: ").strip()
        if user_input in ("終了", "q", "quit"):
            break
        messages.append({"role": "user", "content": user_input})

    return messages


# ─── フェーズ2：スライド構成生成 ─────────────────────

SLIDE_SYSTEM = """\
あなたは高校「公共」「歴史総合」の授業スライドを設計するAIです。
インタビュー内容と提供教材をもとに授業スライドの構成をJSONで出力してください。

【授業設計の方針（厳守）】
- 教師の説明は「生徒が考えるための前提」。説明だけで完結させない
- 冒頭に先行オーガナイザー（単元全体の構造）を organizer スライドで必ず1枚入れる
- question スライドで概念化・抽象化を促す深い問いを設定（一問一答禁止）
- 終盤に activity スライドで「自分の言葉でまとめる」活動を入れる

【出力形式（JSONのみ、前後の説明文不要）】
{
  "title": "授業タイトル",
  "slides": [
    {"type": "title",     "title": "タイトル", "body": "サブタイトル"},
    {"type": "organizer", "title": "先行オーガナイザーのタイトル",
     "bullets": ["学習内容①", "学習内容②", "学習内容③"]},
    {"type": "content",   "title": "スライドタイトル",
     "bullets": ["項目1", "項目2", "項目3"]},
    {"type": "question",  "title": "問い（短く印象的に）", "body": "補足"},
    {"type": "activity",  "title": "活動タイトル", "body": "活動の指示"},
    {"type": "summary",   "title": "まとめ",
     "bullets": ["まとめ1", "まとめ2", "まとめ3"]}
  ]
}

【ルール】
- 合計12〜15枚
- bullets は最大4項目・1項目40字以内
- question を2〜3枚（概念化・抽象化を促す深い問い）
- activity を1〜2枚
- organizer は title の直後に必ず1枚
- JSONのみ出力
"""


def generate_slide_structure(title: str, interview_history: list[dict],
                             textbook: str, memo: str) -> dict:
    interview_text = "\n".join(
        f"{'AI' if m['role'] == 'assistant' else '教師'}: {m['content']}"
        for m in interview_history
    )

    print("\n[スライド構成を生成中...]\n")

    response = client.messages.create(
        model=MODEL,
        max_tokens=4096,
        system=SLIDE_SYSTEM,
        messages=[{"role": "user", "content": (
            f"授業タイトル：{title}\n\n"
            f"【インタビュー内容】\n{interview_text}\n\n"
            f"【教科書テキスト】\n{textbook[:3000] if textbook else '（未提供）'}\n\n"
            f"【NotebookLMまとめ】\n{memo[:3000] if memo else '（未提供）'}\n\n"
            "スライド構成をJSONで出力してください。"
        )}],
    )

    raw = response.content[0].text
    match = re.search(r"\{.*\}", raw, re.DOTALL)
    if not match:
        raise ValueError(f"JSONが見つかりませんでした:\n{raw}")
    return json.loads(match.group())


# ─── メイン ─────────────────────────────────────

def main():
    print("=" * 50)
    print("  授業スライド自動生成ツール（歴史総合・公共）")
    print("=" * 50)

    title = input("\n授業タイトルを入力してください: ").strip()
    if not title:
        print("タイトルが入力されていません。終了します。")
        return

    textbook = load_dir(INPUT_TEXTS)
    memo = load_dir(INPUT_MEMOS)

    if textbook:
        print(f"[教科書テキスト: {len(textbook)}文字]")
    if memo:
        print(f"[NotebookLMまとめ: {len(memo)}文字]")

    interview_history = run_interview(title, textbook, memo)
    structure = generate_slide_structure(title, interview_history, textbook, memo)

    print("\n【生成されたスライド構成】")
    for i, s in enumerate(structure.get("slides", []), 1):
        print(f"  {i:2}. [{s.get('type','content'):10}] {s.get('title','')}")

    confirm = input("\nこの構成でPowerPointを生成しますか？ [Y/n]: ").strip().lower()
    if confirm == "n":
        print("キャンセルしました。")
        return

    timestamp = datetime.now().strftime("%Y%m%d_%H%M")
    safe_title = re.sub(r'[\\/:*?"<>|]', "", title)[:40]
    output_path = OUTPUT_DIR / f"{timestamp}_{safe_title}.pptx"

    create_pptx(structure, output_path)


if __name__ == "__main__":
    main()
