"""
研究資料フォルダ（docx / pdf / pptx）をテキスト抽出してObsidian Vaultの
Markdownノートに変換する。

使い方:
  python import_research_docs.py <フォルダ or ZIPのパス>
  （引数なしの場合は _inbox フォルダを対象にする）

対応形式: .docx, .pdf, .pptx
"""
import re
import sys
import zipfile
import tempfile
from pathlib import Path

import docx
import fitz  # PyMuPDF
from pptx import Presentation

INBOX_DIR = Path(__file__).parent / "_inbox"
VAULT_OUT_DIR = Path.home() / "Desktop" / "MyObsidian" / "研究資料"

SUPPORTED_EXTS = {".docx", ".pdf", ".pptx"}


def sanitize_filename(text: str, maxlen: int = 60) -> str:
    text = re.sub(r'[\\/:*?"<>|\n\r\t]', " ", text or "").strip()
    text = re.sub(r"\s+", " ", text)
    return text[:maxlen] if text else "無題"


def extract_docx(path: Path) -> str:
    document = docx.Document(str(path))
    return "\n\n".join(p.text for p in document.paragraphs if p.text.strip())


def extract_pdf(path: Path) -> str:
    with fitz.open(str(path)) as pdf:
        return "\n\n".join(page.get_text().strip() for page in pdf if page.get_text().strip())


def extract_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    slides_text = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if texts:
            slides_text.append(f"### スライド {i}\n\n" + "\n\n".join(texts))
    return "\n\n".join(slides_text)


EXTRACTORS = {".docx": extract_docx, ".pdf": extract_pdf, ".pptx": extract_pptx}


def convert_file(path: Path, root: Path, out_dir: Path) -> bool:
    try:
        body = EXTRACTORS[path.suffix.lower()](path)
    except Exception as e:
        print(f"  スキップ（読み取り失敗: {e}）: {path.name}")
        return False

    if not body.strip():
        print(f"  スキップ（本文なし）: {path.name}")
        return False

    title = path.stem
    rel_source = path.relative_to(root).as_posix()
    out_path = out_dir / f"{sanitize_filename(title)}.md"

    lines = [
        "---",
        f"title: {title}",
        f"source: {rel_source}",
        "tags: [研究資料]",
        "---",
        "",
        f"# {title}",
        "",
        body,
    ]
    out_path.write_text("\n".join(lines), encoding="utf-8")
    print(f"  変換: {out_path.name}")
    return True


def resolve_source_dir(path: Path) -> Path:
    if path.is_file() and path.suffix.lower() == ".zip":
        tmp_dir = Path(tempfile.mkdtemp(prefix="research_docs_"))
        with zipfile.ZipFile(path) as zf:
            zf.extractall(tmp_dir)
        return tmp_dir
    return path


def main():
    source_arg = Path(sys.argv[1]) if len(sys.argv) > 1 else INBOX_DIR
    if not source_arg.exists():
        if source_arg == INBOX_DIR:
            INBOX_DIR.mkdir(parents=True)
            print(f"{INBOX_DIR} を作成しました。ここにフォルダ/ZIPを置いてから再実行してください。")
        else:
            print(f"{source_arg} が見つかりません。")
        return

    source_dir = resolve_source_dir(source_arg)
    files = [p for p in source_dir.rglob("*") if p.suffix.lower() in SUPPORTED_EXTS]
    if not files:
        print(f"{source_dir} に対応ファイル（.docx/.pdf/.pptx）が見つかりません。")
        return

    VAULT_OUT_DIR.mkdir(parents=True, exist_ok=True)
    print(f"{len(files)}件のファイルを変換します → {VAULT_OUT_DIR}")
    count = 0
    for path in sorted(files):
        print(f"{path.relative_to(source_dir)}")
        if convert_file(path, source_dir, VAULT_OUT_DIR):
            count += 1

    print(f"\n完了: {count}/{len(files)}件を {VAULT_OUT_DIR} に変換しました。")


if __name__ == "__main__":
    main()
